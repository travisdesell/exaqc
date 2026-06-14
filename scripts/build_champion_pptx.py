"""Build a PowerPoint deck of champion-circuit PNGs across every stage.

Output:
  src/Ryan_cookin/results/champion_circuits.pptx

Layout:
  1. Title slide.
  2. Architectural overview slide (what each stage does, what's evolved).
  3. Headline numbers slide (mean / best test_acc per stage).
  4..N. One slide per stage variant (A, B, B_multiseed, C, C_multiseed,
     E_v1, E_v2, F_v1, F_v2, F_v2_multiseed, F_v2_big, G), each showing
     the four dataset champion PNGs in a 2x2 grid with the dataset name
     and the cell's test_acc directly underneath each image.

Uses the same champion-cell selection rule as render_champion_circuits.py
(max test_acc per (stage, dataset), tie-break by min test_loss).
"""
from __future__ import annotations

import csv
import os
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor


RESULTS_DIR = Path("src/Ryan_cookin/results")
PNG_DIR = RESULTS_DIR / "champion_circuits"
OUT_PATH = RESULTS_DIR / "champion_circuits.pptx"

DATASETS = ["iris", "wine", "seeds", "breast_cancer"]


# Order matches the per-stage progression in tables.txt + the new Stage G.
STAGES = [
    ("A",                "Stage A: fixed RY-CNOT ansatz, encoder external",                    "stage_a_trained.csv"),
    ("B",                "Stage B: evolved ansatz, encoder external (seed 0)",                  "stage_b.csv"),
    ("B_multiseed",      "Stage B: 3-seed champion (winners only)",                             "stage_b_multiseed.csv"),
    ("C",                "Stage C: evolved ansatz + evolved N, encoder external (seed 0)",      "stage_c.csv"),
    ("C_multiseed",      "Stage C: 3-seed champion (winners only)",                             "stage_c_multiseed.csv"),
    ("E_v1",             "Stage E v1: encoder enc_ry/rx/rz inside genome, depth-interleaved",   "stage_e.csv"),
    ("E_v2",             "Stage E v2: same as v1, n_genomes 300 (budget bump)",                 "stage_e_v2.csv"),
    ("F_v1",             "Stage F v1: enc_rot + scalar Ising entanglers, depth-interleaved",    "stage_f_v1.csv"),
    ("F_v2",             "Stage F v2: universal enc_rot + enc_xx/yy/zz, depth-interleaved",     "stage_f_v2.csv"),
    ("F_v2_multiseed",   "Stage F v2: 3-seed champion",                                         "stage_f_v2_multiseed.csv"),
    ("F_v2_big",         "Stage F v2 budget bump (n_genomes 300)",                              "stage_f_v2_big.csv"),
    ("G",                "Stage G: ENFORCED encoder block then ansatz block (no interleaving)", "stage_g.csv"),
]


def champion_row(csv_path: Path, dataset: str) -> dict | None:
    if not csv_path.exists():
        return None
    with open(csv_path, newline="") as f:
        rows = [r for r in csv.DictReader(f) if r["dataset"] == dataset]
    if not rows:
        return None
    return sorted(rows, key=lambda r: (-float(r["test_acc"]), float(r["test_loss"])))[0]


def add_title_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    tx = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.3), Inches(2.5))
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "Champion circuits across Stages A through G"
    r.font.size = Pt(40)
    r.font.bold = True
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = ("\nbest test_acc cell per (stage, dataset)\n"
               "rendered from saved per-cell .pt weights")
    r2.font.size = Pt(20)
    r2.font.color.rgb = RGBColor(0x44, 0x44, 0x44)
    p3 = tf.add_paragraph()
    p3.alignment = PP_ALIGN.CENTER
    r3 = p3.add_run()
    r3.text = ("\nRyan_cookin / exaqc -- "
               "branch Quantum-Congressman-Vogt-Cooking")
    r3.font.size = Pt(14)
    r3.font.color.rgb = RGBColor(0x88, 0x88, 0x88)


def add_overview_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # Title
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.7))
    tf = title.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "What each stage evolves"
    r.font.size = Pt(28)
    r.font.bold = True

    overview = [
        ("A",        "trained encoder + FIXED ansatz (RY-CNOT chain), N swept manually"),
        ("B",        "trained encoder + EVOLVED ansatz (CircuitGenome), N swept manually"),
        ("C",        "trained encoder + evolved ansatz + EVOLVED N (grow/shrink_register)"),
        ("D",        "adds linear_proj + reupload_euler encoders, consumed by B/C"),
        ("E v1/v2",  "EVOLVED encoder (enc_ry/rx/rz) + evolved ansatz + evolved N -- interleaved"),
        ("F v1",     "enc_rot universal-1q feat-dep + scalar Ising entanglers -- interleaved"),
        ("F v2",     "enc_rot + enc_xx/yy/zz universal 2q feat-dep -- interleaved (data reuploading)"),
        ("G",        "same vocab as F v2, but ENFORCED encoder block then ansatz block"),
    ]
    text_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(12.3), Inches(6.0))
    tf2 = text_box.text_frame
    tf2.word_wrap = True
    for i, (stage, desc) in enumerate(overview):
        para = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
        para.alignment = PP_ALIGN.LEFT
        run = para.add_run()
        run.text = f"{stage:<10}  {desc}"
        run.font.size = Pt(18)
        run.font.name = "Consolas"
        run.font.color.rgb = RGBColor(0x22, 0x22, 0x22)
        para.space_after = Pt(8)


def add_headline_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.7))
    p = title.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = "Headline numbers (test_acc)"
    r.font.size = Pt(28)
    r.font.bold = True

    rows = [
        ("stage",            "n cells", "seeds", "mean acc", "best", "notes"),
        ("A",                "64",       "1",     "0.707",     "0.967", "baseline, fixed ansatz"),
        ("B all encoders",   "96",       "1",     "0.794",     "1.000", ""),
        ("B winners only",   "144",      "3",     "0.954",     "1.000", "learned 0.903 / linear_proj 0.977 / reupload_euler 0.981"),
        ("C all encoders",   "24",       "1",     "0.862",     "1.000", ""),
        ("C winners only",   "36",       "3",     "0.952",     "1.000", "learned 0.891 / linear_proj 0.976 / reupload_euler 0.989"),
        ("E v1",             "4",        "1",     "0.790",     "0.851", "120 genomes, weak encoder gates"),
        ("E v2",             "4",        "1",     "0.787",     "0.877", "300 genomes -- no improvement"),
        ("F v1",             "4",        "1",     "0.804",     "0.842", "universal-1q only"),
        ("F v2",             "4",        "1",     "0.904",     "0.967", "universal-1q + feat-dep 2q"),
        ("F v2 (3 seeds)",   "12",       "3",     "0.891",     "1.000", "error bars on v2"),
        ("F v2 (300 gen)",   "4",        "1",     "0.899",     "0.967", "breast_cancer +0.035"),
        ("G",                "4",        "1",     "0.873",     "0.921", "block-enforced (no interleaving)"),
    ]
    text_box = slide.shapes.add_textbox(Inches(0.3), Inches(1.0), Inches(12.7), Inches(6.2))
    tf = text_box.text_frame
    tf.word_wrap = False
    for i, row in enumerate(rows):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = para.add_run()
        run.text = f"{row[0]:<22} {row[1]:>10} {row[2]:>7} {row[3]:>12} {row[4]:>7}   {row[5]}"
        run.font.size = Pt(13)
        run.font.name = "Consolas"
        if i == 0:
            run.font.bold = True


def add_stage_slide(prs: Presentation, stage_name: str, stage_desc: str, csv_name: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    title = slide.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(13.0), Inches(0.6))
    p = title.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = stage_desc
    r.font.size = Pt(20)
    r.font.bold = True

    # 2x2 grid of PNG champions. Slide is 13.33" wide x 7.5" tall (16:9).
    # Each image gets roughly half-width / (~3" tall) with room for a caption.
    csv_path = RESULTS_DIR / csv_name
    cell_w = Inches(6.4)
    cell_h = Inches(2.9)
    caption_h = Inches(0.35)
    positions = [
        (Inches(0.2), Inches(0.95)),
        (Inches(6.7), Inches(0.95)),
        (Inches(0.2), Inches(4.15)),
        (Inches(6.7), Inches(4.15)),
    ]
    for (left, top), dataset in zip(positions, DATASETS):
        png_path = PNG_DIR / stage_name / f"{dataset}.png"
        if png_path.exists():
            slide.shapes.add_picture(str(png_path), left, top, width=cell_w, height=cell_h)
        # caption
        row = champion_row(csv_path, dataset)
        if row is None:
            caption = f"{dataset}: (no cell)"
        else:
            acc = float(row["test_acc"])
            n_q = row.get("n_qubits") or row.get("best_n_qubits") or "?"
            enc = row.get("encoder", "")
            enc_str = f"  encoder={enc}" if enc else ""
            caption = f"{dataset}   N={n_q}   test_acc={acc:.3f}{enc_str}"
        cap = slide.shapes.add_textbox(left, top + cell_h, cell_w, caption_h)
        p = cap.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = caption
        r.font.size = Pt(12)
        r.font.name = "Consolas"


def main() -> int:
    prs = Presentation()
    # 16:9 widescreen (13.33" x 7.5")
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    add_title_slide(prs)
    add_overview_slide(prs)
    add_headline_slide(prs)
    for stage_name, stage_desc, csv_name in STAGES:
        add_stage_slide(prs, stage_name, stage_desc, csv_name)

    OUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT_PATH))
    print(f"Wrote {len(prs.slides)} slides to {OUT_PATH}")
    return 0


if __name__ == "__main__":
    sys.exit(main())
