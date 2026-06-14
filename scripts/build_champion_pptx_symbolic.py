"""Build a PowerPoint deck of the SYMBOLIC champion circuits (from the
TikZ PDFs in champion_circuits_tikz/), so the encoder gates show their
data-dependent form (a*x_i + b) and are color-coded blue vs the
orange-coded ansatz gates.

For Stages E/F the encoder and ansatz gates interleave by evolved depth
-- the coloring makes that visible. For Stage G the qnode enforces the
encoder-block-then-ansatz-block ordering, so the colors form clean
contiguous regions. For Stages A/B/C the TikZ only shows the genome's
ansatz portion (orange everywhere), with the external encoder
summarized in the header comment of the source .tex.

Pipeline:
  1. Rasterize each champion_circuits_tikz/<stage>/<dataset>.pdf to PNG
     at 220 DPI using PyMuPDF.
  2. Build a 16:9 .pptx with title + overview + headline + one slide
     per stage variant showing the four dataset champions in 2x2 grid.

Output:
  src/Ryan_cookin/results/champion_circuits_tikz_png/<stage>/<dataset>.png
  src/Ryan_cookin/results/champion_circuits_symbolic.pptx
"""
from __future__ import annotations

import csv
import sys
from pathlib import Path

import fitz  # PyMuPDF
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor


RESULTS_DIR = Path("src/Ryan_cookin/results")
TIKZ_PDF_DIR = RESULTS_DIR / "champion_circuits_tikz"
PNG_DIR = RESULTS_DIR / "champion_circuits_tikz_png"
OUT_PATH = RESULTS_DIR / "champion_circuits_symbolic.pptx"

DATASETS = ["iris", "wine", "seeds", "breast_cancer"]

STAGES = [
    ("A",                "Stage A: Quantikz symbolic (ansatz only; external encoder noted in source .tex)",     "stage_a_trained.csv"),
    ("B",                "Stage B: Quantikz symbolic (ansatz only; external encoder noted in source .tex)",     "stage_b.csv"),
    ("B_multiseed",      "Stage B 3-seed champion (Quantikz)",                                                  "stage_b_multiseed.csv"),
    ("C",                "Stage C: Quantikz symbolic (ansatz only; external encoder noted in source .tex)",     "stage_c.csv"),
    ("C_multiseed",      "Stage C 3-seed champion (Quantikz)",                                                  "stage_c_multiseed.csv"),
    ("E_v1",             "Stage E v1: enc_ry/rx/rz (blue) interleaved with ansatz rotations (orange)",          "stage_e.csv"),
    ("E_v2",             "Stage E v2: same as v1 with n_genomes=300",                                           "stage_e_v2.csv"),
    ("F_v1",             "Stage F v1: enc_rot (blue) + scalar Ising entanglers (orange) -- interleaved",        "stage_f_v1.csv"),
    ("F_v2",             "Stage F v2: enc_rot + enc_xx/yy/zz (blue) + ansatz rotations (orange) -- interleaved","stage_f_v2.csv"),
    ("F_v2_multiseed",   "Stage F v2 3-seed champion",                                                          "stage_f_v2_multiseed.csv"),
    ("F_v2_big",         "Stage F v2 budget bump (n_genomes=300)",                                              "stage_f_v2_big.csv"),
    ("G",                "Stage G: encoder block (blue) THEN ansatz block (orange) -- enforced by qnode",       "stage_g.csv"),
]


def champion_row(csv_path: Path, dataset: str) -> dict | None:
    if not csv_path.exists():
        return None
    with open(csv_path, newline="") as f:
        rows = [r for r in csv.DictReader(f) if r["dataset"] == dataset]
    if not rows:
        return None
    return sorted(rows, key=lambda r: (-float(r["test_acc"]), float(r["test_loss"])))[0]


def rasterize_all_pdfs(dpi: int = 220) -> int:
    """Render each champion_circuits_tikz/<stage>/<dataset>.pdf to PNG."""
    PNG_DIR.mkdir(parents=True, exist_ok=True)
    count = 0
    zoom = dpi / 72.0  # PyMuPDF uses 72 DPI as base
    matrix = fitz.Matrix(zoom, zoom)
    for stage_name, _, _ in STAGES:
        out_stage = PNG_DIR / stage_name
        out_stage.mkdir(exist_ok=True)
        for dataset in DATASETS:
            pdf_path = TIKZ_PDF_DIR / stage_name / f"{dataset}.pdf"
            if not pdf_path.exists():
                print(f"  MISSING  {pdf_path}")
                continue
            png_path = out_stage / f"{dataset}.png"
            doc = fitz.open(str(pdf_path))
            page = doc.load_page(0)
            pix = page.get_pixmap(matrix=matrix, alpha=False)
            pix.save(str(png_path))
            doc.close()
            count += 1
            print(f"  rendered {stage_name:18s} {dataset:14s}  {pix.width}x{pix.height}px")
    return count


def add_title_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tx = slide.shapes.add_textbox(Inches(0.5), Inches(2.3), Inches(12.3), Inches(3.0))
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "Champion circuits: symbolic Quantikz view, Stages A-G"
    r.font.size = Pt(36)
    r.font.bold = True
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = ("\nBLUE boxes = encoding gates (angle = a*x_i + b, data-dependent)"
               "\nORANGE boxes = ansatz / learning gates (trained scalar)"
               "\ncontrol/target idioms (cx, cz) keep their native look"
               "\n\nLook for interleaved blue/orange in Stages E/F (data re-uploading)"
               "\nand clean block separation in Stage G (block-enforced).")
    r2.font.size = Pt(17)
    r2.font.color.rgb = RGBColor(0x33, 0x33, 0x33)


def add_overview_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.7))
    p = title.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = "How encoding and learning gates relate per stage"
    r.font.size = Pt(26)
    r.font.bold = True

    overview = [
        ("A",        "encoder runs in Python (external block) ; ansatz = fixed RY-CNOT chain"),
        ("B",        "encoder in Python (block) ; ansatz = evolved genome"),
        ("C",        "encoder in Python (block) ; ansatz = evolved genome ; evolved N"),
        ("D",        "encoder classes (linear_proj, reupload_euler), consumed by B/C"),
        ("E v1/v2",  "encoder + ansatz gates inside ONE genome, sorted by depth -- INTERLEAVED"),
        ("F v1",     "enc_rot universal-1q (blue) + scalar Ising entanglers (orange) -- INTERLEAVED"),
        ("F v2",     "enc_rot + enc_xx/yy/zz feat-dep (blue) + ansatz (orange) -- INTERLEAVED (data reuploading)"),
        ("G",        "same vocab as F v2, but qnode ENFORCES encoder-pass then ansatz-pass"),
    ]
    text_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(12.3), Inches(6.0))
    tf = text_box.text_frame
    tf.word_wrap = True
    for i, (stage, desc) in enumerate(overview):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = PP_ALIGN.LEFT
        run = para.add_run()
        run.text = f"{stage:<10}  {desc}"
        run.font.size = Pt(16)
        run.font.name = "Consolas"
        run.font.color.rgb = RGBColor(0x22, 0x22, 0x22)
        para.space_after = Pt(8)

    # legend at the bottom
    legend = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.7))
    lf = legend.text_frame
    lp = lf.paragraphs[0]
    lr = lp.add_run()
    lr.text = ("color coding inside the per-stage slides:   "
               "BLUE = encoding gates (data-dependent)   "
               "ORANGE = ansatz / learning gates (trained scalar)")
    lr.font.size = Pt(14)
    lr.font.italic = True
    lr.font.color.rgb = RGBColor(0x55, 0x55, 0x55)


def add_stage_slide(prs: Presentation, stage_name: str, stage_desc: str, csv_name: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title = slide.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(13.0), Inches(0.6))
    p = title.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = stage_desc
    r.font.size = Pt(18)
    r.font.bold = True

    csv_path = RESULTS_DIR / csv_name
    cell_w = Inches(6.4)
    cell_h = Inches(2.85)
    caption_h = Inches(0.4)
    positions = [
        (Inches(0.2), Inches(0.95)),
        (Inches(6.7), Inches(0.95)),
        (Inches(0.2), Inches(4.20)),
        (Inches(6.7), Inches(4.20)),
    ]
    for (left, top), dataset in zip(positions, DATASETS):
        png_path = PNG_DIR / stage_name / f"{dataset}.png"
        if png_path.exists():
            slide.shapes.add_picture(str(png_path), left, top, width=cell_w, height=cell_h)
        row = champion_row(csv_path, dataset)
        if row is None:
            caption = f"{dataset}: (no cell)"
        else:
            acc = float(row["test_acc"])
            n_q = row.get("n_qubits") or row.get("best_n_qubits") or "?"
            enc = row.get("encoder", "")
            enc_str = f"  encoder={enc}" if enc else ""
            caption = f"{dataset}   N={n_q}   test_acc={acc:.3f}{enc_str}"
        cap = slide.shapes.add_textbox(left, top + cell_h, cell_w, caption_h)
        p = cap.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = caption
        r.font.size = Pt(12)
        r.font.name = "Consolas"


def main() -> int:
    print("step 1: rasterize Quantikz PDFs to PNGs")
    n = rasterize_all_pdfs(dpi=220)
    print(f"  rendered {n} PNGs into {PNG_DIR}\n")

    print("step 2: build slides")
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    add_title_slide(prs)
    add_overview_slide(prs)
    for stage_name, stage_desc, csv_name in STAGES:
        add_stage_slide(prs, stage_name, stage_desc, csv_name)

    OUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT_PATH))
    print(f"  wrote {len(prs.slides)} slides to {OUT_PATH}")
    return 0


if __name__ == "__main__":
    sys.exit(main())
